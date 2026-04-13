from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/Users/xiejz959/Xiejz/College/Signal System & Probability/project/SSP_CP")
SLIDE_DIR = ROOT / "Slide"
CHART_DIR = ROOT / "Charts" / "analysis_round1"
OUTPUT = SLIDE_DIR / "Voice_Noise_Presentation_1.1.pptx"

BG_LIGHT = RGBColor(245, 245, 247)
BG_DARK = RGBColor(10, 10, 12)
TEXT_DARK = RGBColor(20, 20, 20)
TEXT_LIGHT = RGBColor(245, 245, 247)
TEXT_MID = RGBColor(104, 104, 109)
ACCENT_BLUE = RGBColor(0, 113, 227)
ACCENT_RED = RGBColor(255, 69, 58)
ACCENT_GREEN = RGBColor(48, 209, 88)
ACCENT_ORANGE = RGBColor(255, 159, 10)


def set_bg(slide, dark: bool = False) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK if dark else BG_LIGHT


def add_title(slide, title: str, subtitle: str | None = None, dark: bool = False) -> None:
    color = TEXT_LIGHT if dark else TEXT_DARK
    subcolor = RGBColor(210, 210, 214) if dark else TEXT_MID

    box = slide.shapes.add_textbox(Inches(0.9), Inches(0.6), Inches(11.2), Inches(1.1))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Avenir Next"
    r.font.bold = True
    r.font.size = Pt(26)
    r.font.color.rgb = color

    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.92), Inches(1.52), Inches(8.6), Inches(0.5))
        tf2 = box2.text_frame
        tf2.clear()
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Avenir Next"
        r2.font.size = Pt(11)
        r2.font.color.rgb = subcolor


def add_big_center_text(slide, title: str, body: str, dark: bool = False) -> None:
    color = TEXT_LIGHT if dark else TEXT_DARK
    subcolor = RGBColor(215, 215, 220) if dark else TEXT_MID
    box = slide.shapes.add_textbox(Inches(0.9), Inches(1.25), Inches(11.2), Inches(2.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Avenir Next"
    r.font.bold = True
    r.font.size = Pt(30)
    r.font.color.rgb = color

    box2 = slide.shapes.add_textbox(Inches(0.95), Inches(3.0), Inches(8.7), Inches(1.4))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = "Avenir Next"
    r2.font.size = Pt(16)
    r2.font.color.rgb = subcolor


def add_bullets(slide, items: list[str], x: float, y: float, w: float, h: float, dark: bool = False, size: int = 18) -> None:
    color = TEXT_LIGHT if dark else TEXT_DARK
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Avenir Next"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)


def add_caption(slide, text: str, x: float, y: float, w: float, dark: bool = False, align=PP_ALIGN.LEFT) -> None:
    color = RGBColor(215, 215, 220) if dark else TEXT_MID
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.45))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Avenir Next"
    r.font.size = Pt(10)
    r.font.color.rgb = color


def add_light_panel(slide, x: float, y: float, w: float, h: float) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(250, 250, 252)
    shape.line.color.rgb = RGBColor(36, 36, 38)


def add_image(slide, path: Path, x: float, y: float, w: float, h: float) -> None:
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def add_audio_placeholder(slide, label: str, filename: str, x: float, y: float, w: float = 3.2, h: float = 0.9, fill_rgb: RGBColor = ACCENT_BLUE) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = fill_rgb
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"{label}\nInsert {filename}"
    r.font.name = "Avenir Next"
    r.font.bold = True
    r.font.size = Pt(13)
    r.font.color.rgb = TEXT_LIGHT


def add_pill(slide, text: str, x: float, y: float, fill_rgb: RGBColor, dark_text: bool = False) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.9), Inches(0.38))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = fill_rgb
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Avenir Next"
    r.font.bold = True
    r.font.size = Pt(10.5)
    r.font.color.rgb = TEXT_DARK if dark_text else TEXT_LIGHT


def add_arrow_label(slide, text: str, x: float, y: float, w: float = 1.75) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Avenir Next"
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = RGBColor(170, 170, 176)


def add_footer(slide, page_no: int, dark: bool = False) -> None:
    color = RGBColor(180, 180, 185) if dark else RGBColor(130, 130, 135)
    box = slide.shapes.add_textbox(Inches(0.9), Inches(6.85), Inches(11.3), Inches(0.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"ECE0402  •  Slide {page_no}"
    r.font.name = "Avenir Next"
    r.font.size = Pt(9)
    r.font.color.rgb = color


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# 1
s = prs.slides.add_slide(blank)
set_bg(s, dark=True)
add_big_center_text(
    s,
    "Designing a Noise Reduction Filter\nfor Voice Signals",
    "A simulated public-announcement case study using frequency-domain analysis and a basic low-pass filter.",
    dark=True,
)
add_caption(s, "ECE0402 • Signals, Systems & Probability", 0.95, 5.55, 6.0, dark=True)
add_caption(s, "Insert team names here", 0.95, 5.95, 4.0, dark=True)
add_footer(s, 1, dark=True)

# 2
s = prs.slides.add_slide(blank)
set_bg(s, dark=True)
add_title(s, "Why this problem matters", "Motivation", dark=True)
add_bullets(
    s,
    [
        "Voice recordings and public announcements often lose clarity when background noise is present.",
        "High-frequency hiss is a useful starting point because it is common in recording chains and easier to analyze.",
        "A simple, explainable filter is enough to show how Signals and Systems ideas work in a real application.",
    ],
    0.95, 1.8, 6.2, 3.6, dark=True, size=19
)
add_pill(s, "Real problem", 8.6, 2.0, ACCENT_BLUE)
add_pill(s, "Controllable model", 8.6, 2.55, ACCENT_ORANGE)
add_pill(s, "Explainable system", 8.6, 3.10, ACCENT_GREEN, dark_text=True)
add_footer(s, 2, dark=True)

# 3
s = prs.slides.add_slide(blank)
set_bg(s, dark=True)
add_title(s, "Problem statement", "Input, output, and goal", dark=True)
add_bullets(
    s,
    [
        "Input: noisy voice corrupted by high-frequency hiss.",
        "Output: a cleaner voice signal after filtering.",
        "Goal: suppress hiss while keeping speech as understandable as possible.",
    ],
    0.95, 1.75, 5.55, 2.5, dark=True, size=19
)
s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.95), Inches(2.18), Inches(1.85), Inches(0.55)).fill.solid()
shape = s.shapes[-1]
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.color.rgb = ACCENT_BLUE
shape.text_frame.paragraphs[0].text = "clean voice"
shape.text_frame.paragraphs[0].font.name = "Avenir Next"
shape.text_frame.paragraphs[0].font.bold = True
shape.text_frame.paragraphs[0].font.size = Pt(12)
shape.text_frame.paragraphs[0].font.color.rgb = TEXT_LIGHT
shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
add_arrow_label(s, "+", 8.95, 2.28, 0.4)
s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.45), Inches(2.18), Inches(1.85), Inches(0.55)).fill.solid()
shape = s.shapes[-1]
shape.fill.fore_color.rgb = ACCENT_ORANGE
shape.line.color.rgb = ACCENT_ORANGE
shape.text_frame.paragraphs[0].text = "hiss noise"
shape.text_frame.paragraphs[0].font.name = "Avenir Next"
shape.text_frame.paragraphs[0].font.bold = True
shape.text_frame.paragraphs[0].font.size = Pt(12)
shape.text_frame.paragraphs[0].font.color.rgb = TEXT_DARK
shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
add_arrow_label(s, "=", 8.95, 3.03, 0.4)
s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(2.93), Inches(1.95), Inches(0.55)).fill.solid()
shape = s.shapes[-1]
shape.fill.fore_color.rgb = ACCENT_RED
shape.line.color.rgb = ACCENT_RED
shape.text_frame.paragraphs[0].text = "noisy voice"
shape.text_frame.paragraphs[0].font.name = "Avenir Next"
shape.text_frame.paragraphs[0].font.bold = True
shape.text_frame.paragraphs[0].font.size = Pt(12)
shape.text_frame.paragraphs[0].font.color.rgb = TEXT_LIGHT
shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
add_arrow_label(s, "→", 8.8, 3.78, 0.4)
s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.45), Inches(3.68), Inches(1.95), Inches(0.55)).fill.solid()
shape = s.shapes[-1]
shape.fill.fore_color.rgb = ACCENT_GREEN
shape.line.color.rgb = ACCENT_GREEN
shape.text_frame.paragraphs[0].text = "filtered voice"
shape.text_frame.paragraphs[0].font.name = "Avenir Next"
shape.text_frame.paragraphs[0].font.bold = True
shape.text_frame.paragraphs[0].font.size = Pt(12)
shape.text_frame.paragraphs[0].font.color.rgb = TEXT_DARK
shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
add_caption(s, "Main design question: where should the filter cut off to reduce noise without making speech sound too dull?", 0.95, 5.55, 10.5, dark=True)
add_footer(s, 3, dark=True)

# 4
s = prs.slides.add_slide(blank)
set_bg(s, dark=True)
add_title(s, "Simulated application scenario", "Why we use a controlled setup", dark=True)
add_bullets(
    s,
    [
        "We use a station/public-announcement background as the story, but we simulate the signal instead of recording a complex real environment.",
        "This keeps the noise characteristics controllable, the experiment repeatable, and the result easier to explain.",
        "The setup still preserves the real-world intuition: speech + recording-chain hiss + system processing.",
    ],
    0.95, 1.75, 6.2, 3.6, dark=True, size=18
)
shape = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(1.9), Inches(4.1), Inches(2.6))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(28, 28, 30)
shape.line.color.rgb = RGBColor(72, 72, 76)
tf = shape.text_frame
tf.word_wrap = True
tf.clear()
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Scene framing\n\nA short public-announcement style voice sample is corrupted by high-frequency hiss and then processed by a simple digital filter."
r.font.name = "Avenir Next"
r.font.size = Pt(18)
r.font.color.rgb = TEXT_LIGHT
add_footer(s, 4, dark=True)

# 5
s = prs.slides.add_slide(blank)
set_bg(s, dark=True)
add_title(s, "Signal characteristics", "Voice and hiss do not occupy frequency space in the same way", dark=True)
add_light_panel(s, 0.8, 1.45, 11.7, 5.15)
add_image(s, CHART_DIR / "spectrum_round1.png", 0.95, 1.65, 11.35, 4.7)
add_caption(s, "The simulated hiss raises high-frequency energy in the noisy voice, creating a clean target for basic low-pass filtering.", 0.98, 6.45, 10.8, dark=True)
add_footer(s, 5, dark=True)

# 6
s = prs.slides.add_slide(blank)
set_bg(s, dark=True)
add_title(s, "Why Fourier Transform helps", "From waveform to frequency-domain evidence", dark=True)
add_bullets(
    s,
    [
        "In the time domain, speech and noise are mixed together and hard to separate directly.",
        "The Fourier Transform shows where energy is concentrated across frequencies.",
        "That frequency view tells us why a low-pass filter is reasonable in this specific hiss scenario.",
    ],
    0.95, 1.8, 5.8, 3.4, dark=True, size=18
)
add_image(s, CHART_DIR / "spectrogram_round1.png", 6.95, 1.45, 5.35, 4.9)
add_caption(s, "High-frequency regions become visibly brighter after noise is added and are reduced again after filtering.", 7.0, 6.45, 5.1, dark=True)
add_footer(s, 6, dark=True)

# 7
s = prs.slides.add_slide(blank)
set_bg(s, dark=True)
add_title(s, "Filter design idea", "Why 3200 Hz becomes the main demo cutoff", dark=True)
add_bullets(
    s,
    [
        "2800 Hz removes more hiss, but speech starts sounding dull.",
        "3600 Hz preserves more speech detail, but leaves more high-frequency noise behind.",
        "3200 Hz gives the best presentation balance between audible improvement and naturalness.",
    ],
    0.95, 1.8, 5.15, 3.3, dark=True, size=17
)
add_light_panel(s, 6.15, 1.45, 6.2, 5.15)
add_image(s, CHART_DIR / "parameter_comparison_round1.png", 6.35, 1.6, 5.95, 4.8)
add_caption(s, "We keep 3600 Hz as a comparison version, but 3200 Hz is easier to defend as the classroom demo choice.", 0.98, 5.9, 10.8, dark=True)
add_footer(s, 7, dark=True)

# 8
s = prs.slides.add_slide(blank)
set_bg(s, dark=True)
add_title(s, "System workflow", "A complete input-output view", dark=True)
add_pill(s, "clean voice", 1.1, 2.9, ACCENT_BLUE)
add_arrow_label(s, "→", 2.95, 2.98, 0.35)
add_pill(s, "hiss noise", 3.35, 2.9, ACCENT_ORANGE, dark_text=True)
add_arrow_label(s, "→", 5.2, 2.98, 0.35)
add_pill(s, "noisy voice", 5.65, 2.9, ACCENT_RED)
add_arrow_label(s, "→", 7.5, 2.98, 0.35)
add_pill(s, "low-pass filter", 7.9, 2.9, ACCENT_GREEN, dark_text=True)
add_arrow_label(s, "→", 9.95, 2.98, 0.35)
add_pill(s, "filtered voice", 10.35, 2.9, ACCENT_BLUE)
add_caption(s, "This is not just a software call: the project includes scenario modeling, frequency analysis, parameter comparison, system design, and evaluation.", 1.08, 4.05, 10.8, dark=True)
add_bullets(
    s,
    [
        "Model a realistic but controllable scenario",
        "Observe the signal in the frequency domain",
        "Design and compare candidate filter parameters",
        "Evaluate the result with audio and figures",
    ],
    1.08, 4.82, 6.5, 1.65, dark=True, size=14
)
add_footer(s, 8, dark=True)

# 9
s = prs.slides.add_slide(blank)
set_bg(s, dark=True)
add_title(s, "Demo setup", "Audio will be inserted manually later", dark=True)
add_audio_placeholder(s, "Clean Voice", "clean_voice.wav", 1.0, 2.0, fill_rgb=ACCENT_BLUE)
add_audio_placeholder(s, "Noisy Voice", "noisy_voice.wav", 1.0, 3.15, fill_rgb=ACCENT_RED)
add_audio_placeholder(s, "Filtered Voice", "filtered_voice.wav", 1.0, 4.30, fill_rgb=ACCENT_GREEN)
add_bullets(
    s,
    [
        "Use the same short sample for all three clips.",
        "Play in the order: clean → noisy → filtered.",
        "Then explain the change with waveform and spectrogram evidence.",
    ],
    6.0, 2.05, 5.2, 2.9, dark=True, size=17
)
add_caption(s, "Leave these placeholders intact for manual audio insertion in PowerPoint.", 1.0, 5.9, 6.0, dark=True)
add_footer(s, 9, dark=True)

# 10
s = prs.slides.add_slide(blank)
set_bg(s, dark=True)
add_title(s, "Results", "Audible improvement supported by visual evidence", dark=True)
add_light_panel(s, 0.82, 1.48, 5.88, 2.75)
add_image(s, CHART_DIR / "demo_pair_round1.png", 0.95, 1.55, 5.6, 2.3)
add_light_panel(s, 6.62, 1.48, 5.82, 4.65)
add_image(s, CHART_DIR / "spectrogram_round1.png", 6.75, 1.55, 5.55, 4.2)
add_caption(s, "Waveform comparison", 2.45, 3.95, 2.8, dark=True, align=PP_ALIGN.CENTER)
add_caption(s, "Spectrogram comparison", 8.25, 5.9, 3.1, dark=True, align=PP_ALIGN.CENTER)
add_bullets(
    s,
    [
        "High-frequency hiss is visibly reduced after filtering.",
        "Speech remains understandable in the processed output.",
        "The improvement is clear, but some high-frequency speech detail is softened.",
    ],
    0.98, 4.45, 5.35, 1.65, dark=True, size=14
)
add_footer(s, 10, dark=True)

# 11
s = prs.slides.add_slide(blank)
set_bg(s, dark=True)
add_title(s, "Discussion and limitations", "Why the result is useful but not perfect", dark=True)
add_bullets(
    s,
    [
        "The low-pass filter works because the simulated hiss is mainly high-frequency.",
        "The same mechanism also removes some useful high-frequency speech detail.",
        "A more realistic station environment would contain broader and more complex noise, which basic low-pass filtering cannot handle as cleanly.",
    ],
    0.95, 1.9, 6.15, 3.35, dark=True, size=17
)
add_pill(s, "Suppression", 8.55, 2.2, ACCENT_GREEN, dark_text=True)
add_pill(s, "Preservation", 8.55, 2.85, ACCENT_BLUE)
add_pill(s, "Tradeoff", 8.55, 3.50, ACCENT_ORANGE, dark_text=True)
add_caption(s, "The project is strongest as a controlled case study, not as a full solution for every real-world noise environment.", 0.98, 5.85, 10.8, dark=True)
add_footer(s, 11, dark=True)

# 12
s = prs.slides.add_slide(blank)
set_bg(s, dark=True)
add_big_center_text(
    s,
    "A simple filter can explain a lot.",
    "We modeled a voice-noise problem, used frequency-domain analysis to guide the design, and showed that a basic low-pass filter can reduce hiss while preserving intelligibility.",
    dark=True,
)
add_caption(s, "Q&A", 0.98, 5.75, 1.0, dark=True)
add_footer(s, 12, dark=True)

OUTPUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUTPUT))
print(f"Saved {OUTPUT}")
