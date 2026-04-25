from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/Users/xiejz959/Xiejz/College/Signal System & Probability/project/SSP_CP")
DEMO_DIR = ROOT / "demo"
OUTPUT = DEMO_DIR / "Voice_Noise_Presentation_3.1.pptx"

BG_DARK = RGBColor(10, 10, 12)
PANEL = RGBColor(28, 28, 30)
PANEL_LINE = RGBColor(72, 72, 76)
TEXT_LIGHT = RGBColor(245, 245, 247)
TEXT_MUTED = RGBColor(190, 190, 196)
TEXT_DIM = RGBColor(138, 138, 146)
ACCENT_BLUE = RGBColor(0, 113, 227)
ACCENT_RED = RGBColor(255, 69, 58)
ACCENT_GREEN = RGBColor(48, 209, 88)
ACCENT_ORANGE = RGBColor(255, 159, 10)
ACCENT_CYAN = RGBColor(90, 200, 250)
ACCENT_GOLD = RGBColor(255, 214, 10)
ACCENT_PURPLE = RGBColor(191, 90, 242)


def set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int,
    color: RGBColor = TEXT_LIGHT,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Avenir Next"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_text(slide, title, 0.85, 0.55, 11.5, 0.85, 27, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.88, 1.38, 9.8, 0.38, 11, color=TEXT_MUTED)


def add_footer(slide, page_no: int) -> None:
    add_text(slide, f"ECE0402  •  Slide {page_no}", 0.9, 6.88, 11.5, 0.22, 9, color=TEXT_DIM, align=PP_ALIGN.RIGHT)


def add_bullets(slide, items: list[str], x: float, y: float, w: float, h: float, size: int = 17) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Avenir Next"
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(8)


def add_panel(slide, x: float, y: float, w: float, h: float, label: str | None = None) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = PANEL_LINE
    if label:
        tf = shape.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = "Avenir Next"
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = TEXT_MUTED


def add_placeholder(slide, label: str, x: float, y: float, w: float, h: float) -> None:
    add_panel(slide, x, y, w, h)
    add_text(slide, label, x + 0.2, y + h / 2 - 0.25, w - 0.4, 0.5, 15, color=TEXT_MUTED, bold=True, align=PP_ALIGN.CENTER)


def add_pill(slide, text: str, x: float, y: float, w: float, color: RGBColor, dark_text: bool = False) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Avenir Next"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = BG_DARK if dark_text else TEXT_LIGHT


def add_arrow(slide, x: float, y: float, text: str = "→") -> None:
    add_text(slide, text, x, y, 0.45, 0.3, 18, color=TEXT_DIM, bold=True, align=PP_ALIGN.CENTER)


def add_equation(slide, text: str, x: float, y: float, w: float = 7.0) -> None:
    add_panel(slide, x, y, w, 0.68)
    add_text(slide, text, x + 0.25, y + 0.17, w - 0.5, 0.35, 16, color=TEXT_LIGHT, bold=True, align=PP_ALIGN.CENTER)


def add_slide(prs: Presentation, title: str, subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, title, subtitle)
    return slide


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 1
s = add_slide(prs, "Designing a Noise Reduction Filter\nfor Voice Signals", "From fixed low-pass filtering to Wiener-style frequency-domain masking")
add_text(s, "ECE0402 • Signals, Systems & Probability", 0.9, 5.45, 6.0, 0.35, 12, color=TEXT_MUTED)
add_text(s, "Junze Xie, Ming Chen, Zifan Xu, Yimin Chen & Wanting Luo", 0.9, 5.86, 7.0, 0.35, 12, color=TEXT_MUTED)
add_footer(s, 1)

# 2
s = add_slide(prs, "Why this problem matters", "Voice denoising as a real signal-processing task")
add_bullets(s, [
    "Voice recordings and public announcements often lose clarity when background noise is present.",
    "High-frequency hiss is a useful starting point because it is common, audible, and frequency-domain friendly.",
    "Our goal is to connect a real application to Signals and Systems ideas through a controllable demo.",
], 0.95, 1.85, 6.6, 3.2, 18)
add_pill(s, "Real problem", 8.5, 2.0, 2.2, ACCENT_BLUE)
add_pill(s, "Controllable model", 8.5, 2.65, 2.2, ACCENT_ORANGE, dark_text=True)
add_pill(s, "Explainable system", 8.5, 3.3, 2.2, ACCENT_GREEN, dark_text=True)
add_footer(s, 2)

# 3
s = add_slide(prs, "Scenario and scope", "Real-world motivation, controlled simulation")
add_bullets(s, [
    "Scenario: a public-announcement style voice signal affected by high-frequency hiss.",
    "We simulate the signal instead of recording a complex station environment.",
    "This keeps the noise controllable, the experiment repeatable, and the result easier to analyze.",
], 0.95, 1.85, 6.25, 3.2, 18)
add_placeholder(s, "PLACEHOLDER\nscenario / setup diagram", 7.65, 1.75, 4.55, 3.2)
add_footer(s, 3)

# 4
s = add_slide(prs, "Problem statement", "Input, output, and design goal")
add_pill(s, "clean voice", 1.0, 2.55, 1.8, ACCENT_BLUE)
add_arrow(s, 3.0, 2.62, "+")
add_pill(s, "hiss noise", 3.55, 2.55, 1.8, ACCENT_ORANGE, dark_text=True)
add_arrow(s, 5.65, 2.62, "=")
add_pill(s, "noisy voice", 6.2, 2.55, 1.8, ACCENT_RED)
add_arrow(s, 8.25, 2.62)
add_pill(s, "denoising system", 8.8, 2.55, 2.35, ACCENT_PURPLE)
add_arrow(s, 11.35, 2.62)
add_pill(s, "enhanced voice", 11.75, 2.55, 1.25, ACCENT_GREEN, dark_text=True)
add_text(s, "Goal: suppress hiss while preserving useful speech information.", 1.0, 4.1, 7.2, 0.45, 18, bold=True)
add_text(s, "This is a tradeoff, not a perfect-recovery problem.", 1.0, 4.65, 7.2, 0.35, 15, color=TEXT_MUTED)
add_footer(s, 4)

# 5
s = add_slide(prs, "Signal model", "A Signals and Systems viewpoint")
add_equation(s, "x_noisy[n] = x_clean[n] + n_hiss[n]", 1.0, 2.0, 6.6)
add_bullets(s, [
    "clean voice and hiss are both signals",
    "the noisy voice is a superposition",
    "the filter is an input-output system",
    "frequency response explains what the system keeps or suppresses",
], 1.0, 3.05, 6.4, 2.4, 17)
add_placeholder(s, "PLACEHOLDER\ninput-output system diagram", 8.0, 2.0, 4.2, 2.8)
add_footer(s, 5)

# 6
s = add_slide(prs, "Visualizing the signals", "Waveform, spectrum, and spectrogram")
add_placeholder(s, "waveform\nplaceholder", 0.95, 1.85, 3.65, 2.65)
add_placeholder(s, "FFT spectrum\nplaceholder", 4.85, 1.85, 3.65, 2.65)
add_placeholder(s, "spectrogram\nplaceholder", 8.75, 1.85, 3.65, 2.65)
add_text(s, "Waveform shows time-domain shape. FFT shows overall frequency energy. Spectrogram shows how energy changes over time and frequency.", 1.0, 5.25, 10.8, 0.65, 17, color=TEXT_MUTED)
add_footer(s, 6)

# 7
s = add_slide(prs, "Frequency-domain observation", "Why hiss suggests a filtering strategy")
add_bullets(s, [
    "Clean voice has structured low-to-mid frequency components.",
    "Hiss noise raises the high-frequency background.",
    "The noisy signal makes the high-frequency region visibly stronger.",
], 0.95, 1.85, 5.25, 2.6, 18)
add_placeholder(s, "PLACEHOLDER\nspectrum / spectrogram evidence", 6.65, 1.55, 5.7, 4.4)
add_footer(s, 7)

# 8
s = add_slide(prs, "Baseline: low-pass filter", "A fixed frequency-domain mask")
add_bullets(s, [
    "Low-pass filtering keeps lower frequencies and attenuates higher frequencies.",
    "It is easy to explain and closely related to frequency response.",
    "In mask language, low-pass is a fixed mask: pass below cutoff, suppress above cutoff.",
], 0.95, 1.85, 6.2, 3.0, 18)
add_placeholder(s, "PLACEHOLDER\nlow-pass response / fixed mask", 7.65, 1.85, 4.55, 3.1)
add_footer(s, 8)

# 9
s = add_slide(prs, "Baseline parameter choice", "Why cutoff is a design decision")
add_bullets(s, [
    "Lower cutoff: stronger noise reduction, but speech becomes duller.",
    "Higher cutoff: more speech detail, but more hiss remains.",
    "The chosen cutoff is a balance between suppression and preservation.",
], 0.95, 1.85, 5.65, 2.7, 18)
add_placeholder(s, "PLACEHOLDER\ncutoff comparison figure", 7.0, 1.55, 5.25, 4.45)
add_footer(s, 9)

# 10
s = add_slide(prs, "Why low-pass is not enough", "The baseline is explainable, but coarse")
add_bullets(s, [
    "Low-pass uses the same attenuation pattern for the entire signal.",
    "It does not know whether a time-frequency region is speech-dominant or noise-dominant.",
    "This can reduce hiss, but it can also remove useful high-frequency speech detail.",
], 0.95, 1.95, 6.2, 3.0, 19)
add_text(s, "This motivates a more flexible frequency-domain mask.", 0.98, 5.25, 8.2, 0.55, 20, color=ACCENT_CYAN, bold=True)
add_footer(s, 10)

# 11
s = add_slide(prs, "Frequency-domain masking", "A more flexible version of the same idea")
add_equation(s, "Y(f) = M(f)X(f)", 0.95, 1.8, 4.7)
add_equation(s, "Y(t, f) = M(t, f)X(t, f)", 0.95, 2.75, 4.7)
add_bullets(s, [
    "M close to 1: preserve this component",
    "M close to 0: suppress this component",
    "Low-pass is a fixed mask; frequency masking can be more adaptive",
], 0.95, 3.9, 5.7, 1.7, 16)
add_placeholder(s, "PLACEHOLDER\nmask heatmap", 7.0, 1.55, 5.25, 4.7)
add_footer(s, 11)

# 12
s = add_slide(prs, "Wiener-style mask", "A principled way to design the mask")
add_equation(s, "M(f) = P_speech(f) / (P_speech(f) + P_noise(f))", 0.95, 1.85, 7.0)
add_bullets(s, [
    "If speech power is stronger, the mask moves closer to 1.",
    "If noise power is stronger, the mask moves closer to 0.",
    "This follows the intuition of Wiener filtering, without claiming an industrial-grade denoising system.",
], 0.95, 2.9, 6.7, 2.5, 17)
add_placeholder(s, "PLACEHOLDER\nWiener gain curve / mask", 8.15, 1.85, 4.05, 3.55)
add_footer(s, 12)

# 13
s = add_slide(prs, "Improved method workflow", "From STFT to masked enhanced voice")
add_pill(s, "noisy voice", 0.95, 3.0, 1.55, ACCENT_RED)
add_arrow(s, 2.65, 3.07)
add_pill(s, "STFT", 3.15, 3.0, 1.2, ACCENT_BLUE)
add_arrow(s, 4.5, 3.07)
add_pill(s, "estimate noise", 5.0, 3.0, 1.9, ACCENT_ORANGE, dark_text=True)
add_arrow(s, 7.05, 3.07)
add_pill(s, "compute mask", 7.55, 3.0, 1.85, ACCENT_PURPLE)
add_arrow(s, 9.55, 3.07)
add_pill(s, "apply mask", 10.05, 3.0, 1.65, ACCENT_CYAN, dark_text=True)
add_arrow(s, 11.85, 3.07)
add_pill(s, "iSTFT", 12.1, 3.0, 0.85, ACCENT_GREEN, dark_text=True)
add_text(s, "The improved method still uses the frequency-domain story, but makes the attenuation pattern more selective.", 1.0, 4.4, 10.8, 0.6, 18, color=TEXT_MUTED)
add_footer(s, 13)

# 14
s = add_slide(prs, "Mask visualization", "The main new figure for the upgraded demo")
add_placeholder(s, "Noisy spectrogram\nplaceholder", 0.85, 1.75, 3.75, 3.9)
add_placeholder(s, "Wiener-style mask\nplaceholder", 4.8, 1.75, 3.75, 3.9)
add_placeholder(s, "Masked spectrogram\nplaceholder", 8.75, 1.75, 3.75, 3.9)
add_text(s, "A good result page should show the logic visually: Noisy -> Mask -> Enhanced.", 0.9, 6.05, 10.8, 0.4, 16, color=TEXT_MUTED)
add_footer(s, 14)

# 15
s = add_slide(prs, "Result comparison", "Baseline versus Wiener-style masking")
add_placeholder(s, "Noisy voice\nspectrogram / audio label", 0.85, 1.7, 3.75, 3.7)
add_placeholder(s, "Low-pass result\nspectrogram / audio label", 4.8, 1.7, 3.75, 3.7)
add_placeholder(s, "Mask result\nspectrogram / audio label", 8.75, 1.7, 3.75, 3.7)
add_text(s, "The goal is not to crown a perfect method, but to explain how a more flexible mask changes the filtering behavior.", 0.9, 5.9, 11.0, 0.55, 16, color=TEXT_MUTED)
add_footer(s, 15)

# 16
s = add_slide(prs, "Demo", "Audio slots to fill after the mask experiment is generated")
add_placeholder(s, "AUDIO SLOT\nclean voice", 1.0, 1.8, 2.5, 1.25)
add_placeholder(s, "AUDIO SLOT\nnoisy voice", 3.9, 1.8, 2.5, 1.25)
add_placeholder(s, "AUDIO SLOT\nlow-pass result", 6.8, 1.8, 2.5, 1.25)
add_placeholder(s, "AUDIO SLOT\nmask result", 9.7, 1.8, 2.5, 1.25)
add_bullets(s, [
    "Play the same speech segment for every version.",
    "Then connect what we hear to the spectrogram and mask figures.",
    "Keep the demo short enough that the audience remembers the contrast.",
], 1.0, 4.0, 8.7, 1.8, 17)
add_footer(s, 16)

# 17
s = add_slide(prs, "Limitations", "What this system still cannot solve")
add_bullets(s, [
    "The noise is simulated and simpler than a real station environment.",
    "Mask quality depends on how well the noise power is estimated.",
    "If speech and noise overlap strongly in the same time-frequency region, separation becomes difficult.",
    "More advanced methods may require adaptive filtering or learning-based denoising.",
], 0.95, 1.85, 8.2, 3.4, 18)
add_placeholder(s, "PLACEHOLDER\nlimitations / extension diagram", 9.6, 2.0, 2.75, 2.7)
add_footer(s, 17)

# 18
s = add_slide(prs, "A simple filter can become a richer system.", "Conclusion")
add_bullets(s, [
    "We modeled voice denoising as a Signals and Systems problem.",
    "Low-pass filtering provides a clear baseline.",
    "Frequency-domain masking extends the same idea in a more flexible way.",
    "A Wiener-style mask gives a principled and visual way to design frequency-dependent attenuation.",
], 0.95, 2.0, 9.2, 3.2, 19)
add_footer(s, 18)

OUTPUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUTPUT)
print(OUTPUT)
